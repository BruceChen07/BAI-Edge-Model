from __future__ import annotations

import csv
import io
import json
import shutil
import subprocess
import tempfile
from pathlib import Path

from docx import Document
from openpyxl import load_workbook
from PIL import Image

try:
    import numpy as np
except ImportError:  # pragma: no cover - optional dependency
    np = None

try:
    import fitz  # PyMuPDF
except ImportError:  # pragma: no cover - optional dependency
    fitz = None

try:
    from pptx import Presentation
except ImportError:  # pragma: no cover - optional dependency
    Presentation = None

try:
    from rapidocr_onnxruntime import RapidOCR
except ImportError:  # pragma: no cover - optional dependency
    RapidOCR = None

try:
    import pytesseract
except ImportError:  # pragma: no cover - optional dependency
    pytesseract = None

from app.services.mineru_service import MinerUService


class ParserService:
    def __init__(self) -> None:
        self._ocr_engine = RapidOCR() if RapidOCR is not None else None
        self.mineru_service = MinerUService()

    def extract_text(
        self,
        file_name: str,
        content: bytes,
        enable_ocr: bool = True,
        file_path: Path | None = None,
    ) -> tuple[str, str]:
        ext = Path(file_name).suffix.lower()

        if ext in {".txt", ".md", ".csv", ".json"}:
            return self._parse_textual(ext, content), "done"
        if ext == ".doc":
            return self._parse_doc(content), "done"
        if ext == ".docx":
            return self._parse_docx(content), "done"
        if ext == ".xlsx":
            return self._parse_xlsx(content), "done"
        if ext == ".pptx":
            return self._parse_pptx(content), "done"
        if ext == ".pdf":
            mineru_result = self._try_mineru(file_path, enable_ocr=enable_ocr)
            if mineru_result:
                return mineru_result, "done"
            return self._parse_pdf(content, enable_ocr=enable_ocr), "done" if enable_ocr else "skipped"
        if ext in {".png", ".jpg", ".jpeg", ".bmp", ".tiff"}:
            mineru_result = self._try_mineru(file_path, enable_ocr=enable_ocr)
            if mineru_result:
                return mineru_result, "done"
            return self._parse_image(file_name, content, enable_ocr), "done" if enable_ocr else "skipped"

        preview = content[:2000].decode("utf-8", errors="ignore")
        return f"Extracted placeholder content for {file_name}.\n\n{preview}", "skipped"

    def _try_mineru(self, file_path: Path | None, enable_ocr: bool) -> str:
        if file_path is None:
            return ""
        if not self.mineru_service.is_available():
            return ""
        if not self.mineru_service.supports_file(file_path):
            return ""
        try:
            return self.mineru_service.parse_file(file_path, enable_ocr=enable_ocr)
        except (RuntimeError, subprocess.CalledProcessError, ValueError):
            return ""

    def _parse_textual(self, ext: str, content: bytes) -> str:
        if ext == ".csv":
            reader = csv.reader(io.StringIO(
                content.decode("utf-8", errors="ignore")))
            return "\n".join(" | ".join(cell.strip() for cell in row) for row in reader)
        if ext == ".json":
            try:
                payload = json.loads(content.decode("utf-8", errors="ignore"))
                return json.dumps(payload, ensure_ascii=False, indent=2)
            except json.JSONDecodeError:
                return content.decode("utf-8", errors="ignore")
        return content.decode("utf-8", errors="ignore")

    def _parse_docx(self, content: bytes) -> str:
        document = Document(io.BytesIO(content))
        blocks: list[str] = []
        for paragraph in document.paragraphs:
            text = paragraph.text.strip()
            if text:
                blocks.append(text)
        for table in document.tables:
            for row in table.rows:
                values = [cell.text.strip()
                          for cell in row.cells if cell.text.strip()]
                if values:
                    blocks.append(" | ".join(values))
        return "\n".join(blocks).strip()

    def _parse_doc(self, content: bytes) -> str:
        converted_docx = self._convert_doc_to_docx(content)
        if converted_docx is not None:
            return self._parse_docx(converted_docx)
        return (
            "Legacy .doc parsing requires LibreOffice (`soffice`) to be available. "
            "Please install LibreOffice or re-save the document as .docx."
        )

    def _parse_xlsx(self, content: bytes) -> str:
        workbook = load_workbook(io.BytesIO(content), data_only=True)
        blocks: list[str] = []
        for sheet in workbook.worksheets:
            blocks.append(f"# Sheet: {sheet.title}")
            for row in sheet.iter_rows(values_only=True):
                values = [str(cell).strip()
                          for cell in row if cell is not None and str(cell).strip()]
                if values:
                    blocks.append(" | ".join(values))
        return "\n".join(blocks).strip()

    def _parse_pptx(self, content: bytes) -> str:
        if Presentation is None:
            return "PPTX parser dependency is not installed."
        presentation = Presentation(io.BytesIO(content))
        blocks: list[str] = []
        for index, slide in enumerate(presentation.slides, start=1):
            blocks.append(f"# Slide {index}")
            for shape in slide.shapes:
                text = getattr(shape, "text", "").strip()
                if text:
                    blocks.append(text)
        return "\n".join(blocks).strip()

    def _parse_pdf(self, content: bytes, enable_ocr: bool) -> str:
        if fitz is None:
            return "PDF parser dependency is not installed."
        document = fitz.open(stream=content, filetype="pdf")
        blocks: list[str] = []
        for index, page in enumerate(document, start=1):
            text = page.get_text("text").strip()
            if text:
                blocks.append(f"# Page {index}\n{text}")
                continue
            if enable_ocr:
                ocr_text = self._ocr_pdf_page(page)
                if ocr_text:
                    blocks.append(f"# Page {index}\n{ocr_text}")
        if not blocks and enable_ocr and not self._has_ocr_engine():
            return (
                "OCR is required for scanned PDF parsing, but no local OCR engine is available. "
                "Install Tesseract OCR or RapidOCR to extract scanned content."
            )
        return "\n\n".join(blocks).strip()

    def _parse_image(self, file_name: str, content: bytes, enable_ocr: bool) -> str:
        if enable_ocr:
            ocr_text = self._ocr_image_bytes(content)
            if ocr_text:
                return ocr_text
            if not self._has_ocr_engine():
                return (
                    f"OCR is required for image parsing, but no local OCR engine is available for {file_name}. "
                    "Install Tesseract OCR or RapidOCR."
                )
            return (
                f"OCR engine was enabled but no text was detected from image: {file_name}"
            )
        return f"Image file uploaded without OCR: {file_name}"

    def _convert_doc_to_docx(self, content: bytes) -> bytes | None:
        soffice = shutil.which("soffice")
        if soffice is None:
            return None

        with tempfile.TemporaryDirectory() as temp_dir:
            temp_path = Path(temp_dir)
            source_path = temp_path / "source.doc"
            source_path.write_bytes(content)

            try:
                subprocess.run(
                    [
                        soffice,
                        "--headless",
                        "--convert-to",
                        "docx",
                        "--outdir",
                        str(temp_path),
                        str(source_path),
                    ],
                    check=True,
                    capture_output=True,
                )
            except subprocess.CalledProcessError:
                return None
            converted_path = temp_path / "source.docx"
            if not converted_path.exists():
                return None
            return converted_path.read_bytes()

    def _ocr_pdf_page(self, page: "fitz.Page") -> str:
        pixmap = page.get_pixmap(matrix=fitz.Matrix(2, 2))
        image_bytes = pixmap.tobytes("png")
        return self._ocr_image_bytes(image_bytes)

    def _has_ocr_engine(self) -> bool:
        if self._ocr_engine is not None and np is not None:
            return True
        return pytesseract is not None and shutil.which("tesseract") is not None

    def _ocr_image_bytes(self, image_bytes: bytes) -> str:
        image = Image.open(io.BytesIO(image_bytes)).convert("RGB")
        if self._ocr_engine is not None and np is not None:
            result, _ = self._ocr_engine(np.array(image))
            if result:
                return "\n".join(item[1] for item in result if len(item) > 1 and item[1]).strip()

        if pytesseract is not None:
            try:
                return pytesseract.image_to_string(image, lang="chi_sim+eng").strip()
            except (RuntimeError, pytesseract.TesseractNotFoundError):
                return ""
        return ""
